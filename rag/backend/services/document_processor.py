import os
import time
from typing import Dict, Any, List
import pandas as pd
import docx
import PyPDF2
from pptx import Presentation
import json
import mimetypes
from pathlib import Path

class DocumentProcessor:
    """Handles processing of various document types and ensures metadata is compatible with ChromaDB"""

    def __init__(self):
        self.supported_extensions = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.doc': self._process_docx,
            '.pptx': self._process_pptx,
            '.csv': self._process_csv,
            '.xlsx': self._process_excel,
            '.xls': self._process_excel,
            '.txt': self._process_text,
            '.json': self._process_json
        }

    async def process_document(self, file_path: str, original_filename: str) -> Dict[str, Any]:
        """Process a document and extract content and metadata"""
        start_time = time.time()
        file_extension = Path(file_path).suffix.lower()

        if file_extension not in self.supported_extensions:
            raise ValueError(f"Unsupported file type: {file_extension}")

        file_size = os.path.getsize(file_path)
        mime_type, _ = mimetypes.guess_type(file_path)

        processor = self.supported_extensions[file_extension]
        content, metadata = processor(file_path)

        # Convert all metadata values to str, int, float, bool, or None
        metadata = self._sanitize_metadata(metadata)

        # Create chunks
        chunks = self._create_chunks(content)

        processing_time = time.time() - start_time

        return {
            "filename": original_filename,
            "file_type": file_extension,
            "content": content,
            "chunks": chunks,
            "chunk_count": len(chunks),
            "metadata": {
                **metadata,
                "file_size": file_size,
                "mime_type": mime_type,
                "processing_time": processing_time
            },
            "processing_time": processing_time
        }

    def _sanitize_metadata(self, metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Convert all values to types supported by ChromaDB (str, int, float, bool, None)"""
        sanitized = {}
        for k, v in metadata.items():
            if isinstance(v, (str, int, float, bool)) or v is None:
                sanitized[k] = v
            elif isinstance(v, list):
                sanitized[k] = ', '.join(map(str, v))
            elif isinstance(v, dict):
                sanitized[k] = json.dumps(v)
            else:
                sanitized[k] = str(v)
        return sanitized

    # -------------------- FILE PROCESSORS -------------------- #

    def _process_pdf(self, file_path: str) -> tuple:
        content = ""
        metadata = {"pages": 0, "type": "pdf"}
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                metadata["pages"] = len(pdf_reader.pages)
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    if text:
                        content += text + "\n"
        except Exception as e:
            content = f"Error processing PDF: {str(e)}"
        return content.strip(), metadata

    def _process_docx(self, file_path: str) -> tuple:
        content = ""
        metadata = {"type": "docx", "paragraphs": 0}
        try:
            doc = docx.Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            content = "\n".join(paragraphs)
            metadata["paragraphs"] = len(paragraphs)
        except Exception as e:
            content = f"Error processing DOCX: {str(e)}"
        return content, metadata

    def _process_pptx(self, file_path: str) -> tuple:
        content = ""
        metadata = {"type": "pptx", "slides": 0}
        try:
            prs = Presentation(file_path)
            metadata["slides"] = len(prs.slides)
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                content += "\n".join(slide_text) + "\n\n"
        except Exception as e:
            content = f"Error processing PPTX: {str(e)}"
        return content.strip(), metadata

    def _process_csv(self, file_path: str) -> tuple:
        try:
            df = pd.read_csv(file_path)
            content = f"CSV Dataset Summary:\nRows: {len(df)}, Columns: {len(df.columns)}\nColumns: {', '.join(df.columns)}\n"
            metadata = {
                "type": "csv",
                "rows": len(df),
                "columns": len(df.columns),
                "column_names": df.columns.tolist(),
                "numeric_columns": df.select_dtypes(include=['number']).columns.tolist(),
                "categorical_columns": df.select_dtypes(include=['object']).columns.tolist()
            }
        except Exception as e:
            content = f"Error processing CSV: {str(e)}"
            metadata = {"type": "csv", "error": str(e)}
        return content, metadata

    def _process_excel(self, file_path: str) -> tuple:
        try:
            excel_file = pd.ExcelFile(file_path)
            sheet_names = excel_file.sheet_names
            content = f"Excel File Summary:\nSheets: {', '.join(sheet_names)}\n"
            metadata = {"type": "excel", "sheets": len(sheet_names), "sheet_names": sheet_names}
        except Exception as e:
            content = f"Error processing Excel: {str(e)}"
            metadata = {"type": "excel", "error": str(e)}
        return content, metadata

    def _process_text(self, file_path: str) -> tuple:
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            metadata = {"type": "text", "characters": len(content), "lines": content.count('\n')+1}
        except Exception as e:
            content = f"Error processing text file: {str(e)}"
            metadata = {"type": "text", "error": str(e)}
        return content, metadata

    def _process_json(self, file_path: str) -> tuple:
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
            content = f"JSON Data Summary:\nType: {type(data).__name__}"
            metadata = {"type": "json", "data_type": type(data).__name__}
        except Exception as e:
            content = f"Error processing JSON: {str(e)}"
            metadata = {"type": "json", "error": str(e)}
        return content, metadata

    # -------------------- CHUNKING -------------------- #

    def _create_chunks(self, content: str, chunk_size: int = 1000, overlap: int = 100) -> List[str]:
        """Split content into overlapping chunks for retrieval"""
        if len(content) <= chunk_size:
            return [content]
        chunks = []
        start = 0
        while start < len(content):
            end = start + chunk_size
            if end < len(content):
                for i in range(end, max(start + chunk_size - 200, start), -1):
                    if content[i] in '.!?':
                        end = i + 1
                        break
            chunk = content[start:end].strip()
            if chunk:
                chunks.append(chunk)
            start = end - overlap
        return chunks
